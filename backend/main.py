import os, io, json, uuid, urllib.request, threading, mimetypes, math, re
from urllib.parse import quote
from typing import Optional, List
from fastapi import FastAPI, UploadFile, File, Form, Depends, Request
from fastapi.responses import JSONResponse, StreamingResponse, Response
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from fastembed import TextEmbedding, SparseTextEmbedding
from qdrant_client import QdrantClient
from qdrant_client.models import (Distance, VectorParams, PointStruct, Filter, FieldCondition,
                                  MatchValue, MatchAny, SparseVector, SparseVectorParams, Modifier,
                                  Prefetch, FusionQuery, Fusion)
import auth

QDRANT    = os.environ.get("QDRANT_URL", "http://qdrant:6333")
COLL      = os.environ.get("COLLECTION", "rag_hybrid")
CACHE     = os.environ.get("FASTEMBED_CACHE", "/models")
PREFERRED = os.environ.get("EMBED_MODEL", "BAAI/bge-m3")
FALLBACK  = "intfloat/multilingual-e5-large"
SPARSE_MODEL   = os.environ.get("SPARSE_MODEL", "Qdrant/bm25")
DEFAULT_RERANK = os.environ.get("RERANK_MODEL", "jinaai/jina-reranker-v2-base-multilingual")
RERANK_FALLBACKS = ["jinaai/jina-reranker-v2-base-multilingual", "BAAI/bge-reranker-v2-m3",
                    "Xenova/ms-marco-MiniLM-L-12-v2", "Xenova/ms-marco-MiniLM-L-6-v2"]
CONFIG_PATH = os.environ.get("CONFIG_PATH", "/srv/config.json")
IMG_EXTS = (".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif")
EXTS = (".md", ".txt", ".pdf", ".docx", ".xlsx", ".xls", ".csv", ".pptx") + IMG_EXTS
OCR_LANG = os.environ.get("OCR_LANG", "deu+eng")

DEFAULTS = {
    "default_engine": "local",
    "retrieval": {"top_k": 4, "hybrid": True, "rerank": True, "candidates": 20,
                  "rerank_model": DEFAULT_RERANK},
    "engines": {
        "local":  {"enabled": True,  "ollama_url": os.environ.get("OLLAMA_URL", "http://192.168.1.193:11434"),
                   "model": os.environ.get("OLLAMA_MODEL", "qwen2.5:7b")},
        "claude": {"enabled": False, "api_key": "", "model": "claude-opus-4-8"},
        "openai": {"enabled": False, "api_key": "", "model": "gpt-4o"},
        "gemini": {"enabled": False, "api_key": "", "model": "gemini-2.0-flash"},
    },
    "sources": {"smb": {"enabled": False, "host": "192.168.1.5", "share": "", "path": "", "username": "", "password": ""}},
}
ENV_KEYS = {"claude": "ANTHROPIC_API_KEY", "openai": "OPENAI_API_KEY", "gemini": "GOOGLE_API_KEY"}
LABELS   = {"local": "Lokal · Mac", "claude": "Claude (Anthropic)", "openai": "ChatGPT (OpenAI)", "gemini": "Google Gemini"}
ORDER    = ["local", "claude", "openai", "gemini"]
_lock = threading.Lock()
STATE = {"indexing": False, "detail": "", "result": None, "error": None}
_state_lock = threading.Lock()
def set_indexing(on, detail=""):
    with _state_lock:
        STATE["indexing"] = bool(on); STATE["detail"] = detail

SYSTEM = ("Du bist ein praeziser Wissens-Assistent fuer Tech-IT Consulting. Beantworte die Frage "
          "AUSSCHLIESSLICH anhand des bereitgestellten Kontexts auf Deutsch. Wenn die Antwort nicht "
          "im Kontext steht, sage das ehrlich. Erfinde nichts.")

def _clone(d): return json.loads(json.dumps(d))
def deep_merge(base, ov):
    out = _clone(base)
    for k, v in ov.items():
        out[k] = deep_merge(out[k], v) if isinstance(v, dict) and isinstance(out.get(k), dict) else v
    return out
def load_raw():
    cfg = _clone(DEFAULTS)
    if os.path.exists(CONFIG_PATH):
        try:
            with open(CONFIG_PATH) as f: cfg = deep_merge(cfg, json.load(f))
        except Exception: pass
    return cfg
def load_config():
    cfg = load_raw()
    for e, env in ENV_KEYS.items():
        if not cfg["engines"][e].get("api_key") and os.environ.get(env):
            cfg["engines"][e]["api_key"] = os.environ[env]
    return cfg
def save_config(cfg):
    with _lock:
        with open(CONFIG_PATH, "w") as f: json.dump(cfg, f, indent=2, ensure_ascii=False)
def engine_key(cfg, e): return cfg["engines"][e].get("api_key") or os.environ.get(ENV_KEYS.get(e, ""), "")
def engine_available(cfg, e):
    en = cfg["engines"][e]
    return False if not en.get("enabled") else (True if e == "local" else bool(engine_key(cfg, e)))

def pick_model():
    s = {m["model"] for m in TextEmbedding.list_supported_models()}
    return PREFERRED if PREFERRED in s else FALLBACK
MODEL = pick_model()
EMB   = TextEmbedding(model_name=MODEL, cache_dir=CACHE)
IS_E5 = "e5" in MODEL.lower()
EMBED_BATCH = int(os.environ.get("EMBED_BATCH", "16"))
# parallel=1: kein Multiprocessing -> nur EIN Modell im RAM (sonst lädt fastembed je Worker
# eine eigene ~2 GB-Kopie -> OOM auf der 8-GB-VM). batch_size begrenzt den Spitzenverbrauch.
def emb_query(t):
    x = f"query: {t}" if IS_E5 else t
    return list(EMB.embed([x], parallel=1))[0].tolist()
def emb_passages(texts):
    xs = [f"passage: {t}" for t in texts] if IS_E5 else list(texts)
    return [v.tolist() for v in EMB.embed(xs, batch_size=EMBED_BATCH, parallel=1)]
DIM = len(emb_query("test"))

# ---- Sparse-/BM25-Embeddings (exakte Stichwort-Treffer fuer die Hybrid-Suche) ----
try:
    SEMB = SparseTextEmbedding(model_name=SPARSE_MODEL, cache_dir=CACHE, language="german")
except TypeError:
    SEMB = SparseTextEmbedding(model_name=SPARSE_MODEL, cache_dir=CACHE)
def _to_sparse(e):
    idx = e.indices.tolist() if hasattr(e.indices, "tolist") else list(e.indices)
    val = e.values.tolist()  if hasattr(e.values, "tolist")  else list(e.values)
    return SparseVector(indices=idx, values=val)
def emb_sparse_query(t):    return _to_sparse(list(SEMB.query_embed(t))[0])
def emb_sparse_passages(ts): return [_to_sparse(e) for e in SEMB.embed(list(ts))]

# ---- Reranker (Cross-Encoder, lazy geladen, mehrsprachig) ----
_reranker = None; _reranker_req = None
def get_reranker(name):
    global _reranker, _reranker_req
    if _reranker is not None and _reranker_req == name: return _reranker
    from fastembed.rerank.cross_encoder import TextCrossEncoder
    supported = {m["model"] for m in TextCrossEncoder.list_supported_models()}
    chosen = name if name in supported else next((m for m in RERANK_FALLBACKS if m in supported), None)
    if not chosen: raise RuntimeError("kein unterstuetztes Reranker-Modell verfuegbar")
    _reranker = TextCrossEncoder(model_name=chosen, cache_dir=CACHE); _reranker_req = name
    if chosen != name: print(f"[Rerank] '{name}' nicht verfuegbar -> '{chosen}'")
    return _reranker

qc = QdrantClient(url=QDRANT)
def _create_coll():
    qc.create_collection(COLL,
        vectors_config={"dense": VectorParams(size=DIM, distance=Distance.COSINE)},
        sparse_vectors_config={"bm25": SparseVectorParams(modifier=Modifier.IDF)})
    # Keyword-Index fuer den ACL-Filter (schnelles MatchAny auf den Ordner-Praefixen)
    try: qc.create_payload_index(COLL, "path_prefixes", "keyword")
    except Exception as ex: print(f"[Index-Warn] path_prefixes: {ex}")
def ensure_collection():
    if not qc.collection_exists(COLL): _create_coll()

def src_allowed(user, src):
    """Darf der Nutzer diese Quelle sehen/oeffnen? Nutzt die schon berechneten user['folders']."""
    if not user or user.get("admin"):
        return True
    src = (src or "").replace("\\", "/").strip("/")
    for p in user.get("folders", []):
        if src == p or src.startswith(p + "/"):
            return True
    return False

def acl_filter(user):
    """Qdrant-Filter aus den erlaubten Ordner-Praefixen des Nutzers. Admin -> None (alles)."""
    if not user or user.get("admin"):
        return None
    prefixes = user.get("folders", [])
    if not prefixes:
        # Kein Ordner freigegeben -> bewusst leeres Ergebnis erzwingen
        return Filter(must=[FieldCondition(key="path_prefixes",
                                           match=MatchValue(value="__none__"))])
    return Filter(must=[FieldCondition(key="path_prefixes", match=MatchAny(any=prefixes))])

def _candidates(cfg, q, n, flt=None):
    """Kandidaten holen: hybrid (dense + BM25, RRF-Fusion) oder dense-only als Fallback.
    flt = optionaler ACL-Filter (Ordner-Praefixe des Nutzers)."""
    if cfg.get("retrieval", {}).get("hybrid", True):
        try:
            return qc.query_points(COLL, with_payload=True, limit=n, query_filter=flt,
                prefetch=[Prefetch(query=emb_query(q),        using="dense", limit=n, filter=flt),
                          Prefetch(query=emb_sparse_query(q), using="bm25",  limit=n, filter=flt)],
                query=FusionQuery(fusion=Fusion.RRF)).points
        except Exception as ex:
            print(f"[Hybrid-Warn] Fallback auf Dense-Suche: {ex}")
    return qc.query_points(COLL, query=emb_query(q), using="dense", limit=n,
                           with_payload=True, query_filter=flt).points

def rerank_hits(cfg, q, hits):
    """Cross-Encoder ueber die Kandidaten; Score -> Sigmoid (0..1), absteigend sortiert."""
    rr = get_reranker(cfg.get("retrieval", {}).get("rerank_model", DEFAULT_RERANK))
    scores = list(rr.rerank(q, [h.payload.get("text", "") for h in hits]))
    out = []
    for i in sorted(range(len(hits)), key=lambda j: scores[j], reverse=True):
        h = hits[i]
        try: h.score = 1.0 / (1.0 + math.exp(-float(scores[i])))
        except Exception: pass
        out.append(h)
    return out

def retrieve(cfg, q, user=None):
    rcfg = cfg.get("retrieval", {})
    k = int(rcfg.get("top_k", 4))
    do_rerank = bool(rcfg.get("rerank", True))
    pool = max(int(rcfg.get("candidates", 20)), k) if do_rerank else k
    hits = _candidates(cfg, q, pool, flt=acl_filter(user))
    if hits and do_rerank:
        try: return rerank_hits(cfg, q, hits)[:k]
        except Exception as ex: print(f"[Rerank-Warn] uebersprungen: {ex}")
    return hits[:k]

SOURCE_REL_MARGIN = float(os.environ.get("SOURCE_REL_MARGIN", "0.15"))
def rank_sources(hits):
    """Quellen nach bestem Chunk-Score, nur die relevanten (nahe am Top-Score), min. 1.
       Relativer Abstand -> skalenunabhaengig (passt fuer Rerank-, RRF- und Cosine-Scores)."""
    best = {}  # source -> (score, page) des besten Chunks
    for h in hits:
        src = h.payload.get("source", "?"); sc = float(h.score)
        if src not in best or sc > best[src][0]: best[src] = (sc, h.payload.get("page"))
    ranked = sorted(best.items(), key=lambda x: -x[1][0])
    if not ranked: return []
    thr = ranked[0][1][0] * (1.0 - SOURCE_REL_MARGIN)
    out = []
    for s, (sc, pg) in ranked:
        if sc < thr: continue
        item = {"source": s, "score": round(sc, 3)}
        if pg: item["page"] = pg
        out.append(item)
    return out

def ocr_image_bytes(data):
    import pytesseract
    from PIL import Image
    return pytesseract.image_to_string(Image.open(io.BytesIO(data)), lang=OCR_LANG)

def ocr_pdf_bytes(data, dpi=200):
    """Gescanntes PDF: jede Seite rendern und per Tesseract OCR auslesen."""
    import fitz  # PyMuPDF
    import pytesseract
    from PIL import Image
    parts = []
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        for page in doc:
            pix = page.get_pixmap(dpi=dpi)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            parts.append(pytesseract.image_to_string(img, lang=OCR_LANG))
    finally:
        doc.close()
    return "\n".join(parts)

def ocr_pdf_pages(data, dpi=200):
    """Gescanntes PDF: jede Seite rendern + OCR, mit Seitenzahl."""
    import fitz  # PyMuPDF
    import pytesseract
    from PIL import Image
    out = []
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        for i, page in enumerate(doc, 1):
            pix = page.get_pixmap(dpi=dpi)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            out.append((i, pytesseract.image_to_string(img, lang=OCR_LANG)))
    finally:
        doc.close()
    return out

def read_pages(name, data):
    """[(page, text), ...] – PDFs mit Seitenzahl (OCR-Fallback pro Seite), sonst [(None, Gesamttext)]."""
    ext = ("." + name.lower().rsplit(".", 1)[-1]) if "." in name else ""
    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            pages = [(i, (p.extract_text() or "")) for i, p in enumerate(PdfReader(io.BytesIO(data)).pages, 1)]
            if sum(len(t.strip()) for _, t in pages) < 50:  # kaum Text -> vermutlich gescannt -> OCR pro Seite
                try: pages = ocr_pdf_pages(data)
                except Exception as ex: print(f"[OCR-Warn] {name}: {ex}")
            return pages
        except Exception as ex:
            print(f"[Warn] {name}: {ex}"); return []
    return [(None, read_text_bytes(name, data))]

def read_text_bytes(name, data):
    ext = ("." + name.lower().rsplit(".", 1)[-1]) if "." in name else ""
    try:
        if ext in (".md", ".txt"): return data.decode("utf-8", "ignore")
        if ext == ".pdf":
            from pypdf import PdfReader
            txt = "\n".join((p.extract_text() or "") for p in PdfReader(io.BytesIO(data)).pages)
            if len(txt.strip()) < 50:  # kaum Text -> vermutlich gescannt -> OCR
                try: txt = ocr_pdf_bytes(data)
                except Exception as ex: print(f"[OCR-Warn] {name}: {ex}")
            return txt
        if ext == ".docx":
            import docx
            d = docx.Document(io.BytesIO(data)); parts = [p.text for p in d.paragraphs]
            for t in d.tables:
                for row in t.rows: parts.append(" | ".join(c.text for c in row.cells))
            return "\n".join(parts)
        if ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"# {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c not in (None, "")]
                    if cells: parts.append(" | ".join(cells))
            return "\n".join(parts)
        if ext == ".xls":
            import xlrd
            book = xlrd.open_workbook(file_contents=data); parts = []
            for sh in book.sheets():
                parts.append(f"# {sh.name}")
                for r in range(sh.nrows):
                    cells = [str(sh.cell_value(r, c)) for c in range(sh.ncols) if sh.cell_value(r, c) not in ("", None)]
                    if cells: parts.append(" | ".join(cells))
            return "\n".join(parts)
        if ext == ".csv":
            import csv as _csv
            text = data.decode("utf-8", "ignore")
            try: delim = _csv.Sniffer().sniff(text[:2048], delimiters=",;\t|").delimiter
            except Exception: delim = ";" if text[:2048].count(";") >= text[:2048].count(",") else ","
            rows = _csv.reader(io.StringIO(text), delimiter=delim)
            return "\n".join(" | ".join(c for c in row if c) for row in rows if any(row))
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data)); parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"# Folie {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            t = "".join(run.text for run in p.runs)
                            if t.strip(): parts.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            parts.append(" | ".join(c.text for c in row.cells))
            return "\n".join(parts)
        if ext in IMG_EXTS:
            try: return ocr_image_bytes(data)
            except Exception as ex: print(f"[OCR-Warn] {name}: {ex}")
    except Exception as ex: print(f"[Warn] {name}: {ex}")
    return ""

def chunk(text, size=1000, overlap=150):
    text = text.strip(); out, i = [], 0
    while i < len(text):
        piece = text[i:i+size].strip()
        if piece: out.append(piece)
        i += size - overlap
    return out

def delete_source(rel):
    qc.delete(COLL, points_selector=Filter(must=[FieldCondition(key="source", match=MatchValue(value=rel))]))

def path_prefixes(rel):
    """Alle Vorgaenger-Verzeichnisse von rel (ohne Dateiname) – fuer die Ordner-ACL.
    'Technik/Maschinen/cforce.pdf' -> ['Technik', 'Technik/Maschinen']."""
    parts = rel.replace("\\", "/").strip("/").split("/")[:-1]
    return ["/".join(parts[:i + 1]) for i in range(len(parts))]

def index_file(rel, data, sig=None):
    ensure_collection()
    delete_source(rel)
    items = [(ch, pg) for pg, text in read_pages(rel, data) for ch in chunk(text)]
    if not items: return 0
    texts = [t for t, _ in items]
    tenant = rel.split("/")[0] if "/" in rel else "qnap"
    prefixes = path_prefixes(rel)
    pts = [PointStruct(id=str(uuid.uuid4()), vector={"dense": dv, "bm25": sv},
                       payload={"text": t, "source": rel, "page": pg, "tenant": tenant,
                                "path_prefixes": prefixes, "sig": sig})
           for (t, pg), dv, sv in zip(items, emb_passages(texts), emb_sparse_passages(texts))]
    for i in range(0, len(pts), 256): qc.upsert(COLL, points=pts[i:i+256])
    return len(pts)

def indexed_sigs():
    """Liefert {source: sig} der bereits indizierten Dateien (Signatur = groesse-mtime)."""
    ensure_collection(); out = {}; off = None
    try:
        while True:
            pts, off = qc.scroll(COLL, limit=256, offset=off, with_payload=["source", "sig"])
            for p in pts:
                src = p.payload.get("source")
                if src and src not in out: out[src] = p.payload.get("sig")
            if off is None: break
    except Exception: pass
    return out

def smb_connect(s):
    import smbclient
    smbclient.reset_connection_cache()
    smbclient.register_session(s["host"].strip(), username=s["username"], password=s.get("password", ""),
                               connection_timeout=20)
    host = s["host"].strip(); share = s["share"].strip().strip("\\/")
    sub = (s.get("path") or "").strip().strip("\\/").replace("/", "\\")
    root = f"\\\\{host}\\{share}"
    base = root + ("\\" + sub if sub else "")
    return smbclient, root, base, sub

def file_sig(stat):
    """Signatur aus Groesse + Aenderungszeit – erkennt geaenderte Dateien ohne Download."""
    return f"{getattr(stat, 'st_size', 0)}-{int(getattr(stat, 'st_mtime', 0))}"

def smb_read_retry(s, full_p):
    """Datei lesen; bei abgelaufener SMB-Session (z. B. weil das Embedding zwischen zwei
    Downloads mehrere Minuten dauert) einmal neu verbinden und erneut versuchen."""
    import smbclient
    try:
        with smbclient.open_file(full_p, mode="rb", share_access="r") as fd: return fd.read()
    except Exception:
        smb_connect(s)  # Session erneuern
        with smbclient.open_file(full_p, mode="rb", share_access="r") as fd: return fd.read()

def ingest_smb(full=False):
    cfg = load_config(); s = cfg["sources"]["smb"]
    if not s.get("enabled"): raise ValueError("SMB-Quelle ist nicht aktiviert.")
    if not s.get("share") or not s.get("username"): raise ValueError("Share-Name und Benutzername muessen gesetzt sein.")
    smbclient, root, base, sub = smb_connect(s)

    # 1) Aktuellen Dateibestand auf dem Share inkl. Signatur erfassen (ohne Inhalt zu laden)
    set_indexing(True, "Dateiliste wird gelesen…")
    current = {}  # rel -> sig
    for dp, _d, files in smbclient.walk(base):
        for fn in files:
            if fn.lower().endswith(EXTS):
                full_p = dp + "\\" + fn
                rel = full_p[len(root):].strip("\\").replace("\\", "/")
                try: current[rel] = file_sig(smbclient.stat(full_p))
                except Exception: current[rel] = None

    # 2) Voll-Reindex: Collection neu aufbauen. Sonst: bestehende Signaturen vergleichen
    if full:
        if qc.collection_exists(COLL): qc.delete_collection(COLL)
        _create_coll()
        existing = {}
    else:
        existing = indexed_sigs()

    # 3) Entfernte Dateien aus dem Index loeschen
    removed = [rel for rel in existing if rel not in current]
    for rel in removed: delete_source(rel)

    # 4) Nur geaenderte/neue Dateien herunterladen und neu indizieren
    to_index = [rel for rel, sig in current.items() if full or existing.get(rel) != sig]
    unchanged = len(current) - len(to_index)
    added = updated = skipped = chunks = 0
    for i, rel in enumerate(to_index, 1):
        set_indexing(True, f"{i}/{len(to_index)}: {rel}")
        full_p = root + "\\" + rel.replace("/", "\\")
        try:
            data = smb_read_retry(s, full_p)
        except Exception as ex:
            print(f"[Warn] {rel}: {ex}"); skipped += 1; continue
        was_known = rel in existing
        n = index_file(rel, data, current.get(rel))
        chunks += n
        if not n: skipped += 1
        elif was_known: updated += 1
        else: added += 1

    return {"files": len(current), "added": added, "updated": updated, "unchanged": unchanged,
            "removed": len(removed), "skipped": skipped, "chunks": chunks}

def indexed_count():
    try: return qc.count(COLL).count
    except Exception: return 0

def list_documents(user=None):
    ensure_collection(); srcs = {}; off = None
    try:
        while True:
            pts, off = qc.scroll(COLL, limit=256, offset=off, with_payload=["source"])
            for p in pts:
                src = p.payload.get("source", "?")
                if user is not None and not src_allowed(user, src): continue
                srcs[src] = srcs.get(src, 0) + 1
            if off is None: break
    except Exception: pass
    return [{"source": k, "chunks": srcs[k]} for k in sorted(srcs)]

def all_prefixes():
    """Distinct Ordner-Praefixe ueber den ganzen Index – fuer die Ordnerzuweisung im Admin."""
    ensure_collection(); out = set(); off = None
    try:
        while True:
            pts, off = qc.scroll(COLL, limit=256, offset=off, with_payload=["path_prefixes"])
            for p in pts:
                for pre in (p.payload.get("path_prefixes") or []): out.add(pre)
            if off is None: break
    except Exception: pass
    return sorted(out)

def gen_local(cfg, q, ctx):
    en = cfg["engines"]["local"]
    body = json.dumps({"model": en["model"], "prompt": f"Kontext:\n{ctx}\n\nFrage: {q}\n\nAntwort:", "system": SYSTEM, "stream": False}).encode()
    req = urllib.request.Request(en["ollama_url"] + "/api/generate", data=body, headers={"Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=180) as r: return json.load(r)["response"].strip()
def gen_claude(cfg, q, ctx):
    import anthropic
    en = cfg["engines"]["claude"]; client = anthropic.Anthropic(api_key=engine_key(cfg, "claude"))
    msg = client.messages.create(model=en["model"], max_tokens=1024, system=SYSTEM,
            messages=[{"role": "user", "content": f"Kontext:\n{ctx}\n\nFrage: {q}"}])
    return "".join(b.text for b in msg.content if b.type == "text").strip()
def gen_openai(cfg, q, ctx):
    from openai import OpenAI
    en = cfg["engines"]["openai"]; client = OpenAI(api_key=engine_key(cfg, "openai"))
    r = client.chat.completions.create(model=en["model"], max_tokens=1024,
            messages=[{"role": "system", "content": SYSTEM}, {"role": "user", "content": f"Kontext:\n{ctx}\n\nFrage: {q}"}])
    return (r.choices[0].message.content or "").strip()
def gen_gemini(cfg, q, ctx):
    from google import genai
    en = cfg["engines"]["gemini"]; client = genai.Client(api_key=engine_key(cfg, "gemini"))
    return (client.models.generate_content(model=en["model"], contents=f"{SYSTEM}\n\nKontext:\n{ctx}\n\nFrage: {q}").text or "").strip()
GENERATORS = {"local": gen_local, "claude": gen_claude, "openai": gen_openai, "gemini": gen_gemini}

# ---- Multi-Turn: Verlauf in Chat-Nachrichten umwandeln ----
HISTORY_TURNS = 8  # wie viele vorherige Nachrichten in den Kontext genommen werden
def build_messages(history, q, ctx):
    msgs = [{"role": "system", "content": SYSTEM}]
    for h in (history or [])[-HISTORY_TURNS:]:
        content = (h.get("content") or "").strip()
        role = h.get("role")
        if not content: continue
        if role == "user": msgs.append({"role": "user", "content": content})
        elif role in ("assistant", "bot"): msgs.append({"role": "assistant", "content": content})
    msgs.append({"role": "user", "content": f"Kontext:\n{ctx}\n\nFrage: {q}"})
    return msgs

def _flatten(messages):
    lines = []
    for m in messages:
        if m["role"] == "system": lines.append(m["content"])
        elif m["role"] == "user": lines.append(f"Frage/Nutzer:\n{m['content']}")
        else: lines.append(f"Assistent:\n{m['content']}")
    return "\n\n".join(lines)

# ---- Streaming-Generatoren: liefern Text-Stuecke (Deltas) ----
def stream_local(cfg, messages):
    en = cfg["engines"]["local"]
    body = json.dumps({"model": en["model"], "messages": messages, "stream": True}).encode()
    req = urllib.request.Request(en["ollama_url"] + "/api/chat", data=body, headers={"Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=300) as r:
        for line in r:
            line = line.strip()
            if not line: continue
            try: obj = json.loads(line)
            except Exception: continue
            piece = (obj.get("message") or {}).get("content", "")
            if piece: yield piece
            if obj.get("done"): break
def stream_claude(cfg, messages):
    import anthropic
    en = cfg["engines"]["claude"]; client = anthropic.Anthropic(api_key=engine_key(cfg, "claude"))
    sys = messages[0]["content"] if messages and messages[0]["role"] == "system" else SYSTEM
    msgs = [m for m in messages if m["role"] != "system"]
    with client.messages.stream(model=en["model"], max_tokens=1024, system=sys, messages=msgs) as s:
        for text in s.text_stream: yield text
def stream_openai(cfg, messages):
    from openai import OpenAI
    en = cfg["engines"]["openai"]; client = OpenAI(api_key=engine_key(cfg, "openai"))
    s = client.chat.completions.create(model=en["model"], messages=messages, stream=True)
    for ch in s:
        d = ch.choices[0].delta.content if ch.choices else None
        if d: yield d
def stream_gemini(cfg, messages):
    from google import genai
    en = cfg["engines"]["gemini"]; client = genai.Client(api_key=engine_key(cfg, "gemini"))
    for ch in client.models.generate_content_stream(model=en["model"], contents=_flatten(messages)):
        if getattr(ch, "text", None): yield ch.text
STREAMERS = {"local": stream_local, "claude": stream_claude, "openai": stream_openai, "gemini": stream_gemini}

def sse(event, data):
    return f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"

app = FastAPI(title="Tech-IT Wissens-Chat")
class AskReq(BaseModel):
    question: str
    engine: Optional[str] = None
    history: Optional[List[dict]] = None
class ConfigIn(BaseModel):
    default_engine: Optional[str] = None
    retrieval: Optional[dict] = None
    engines: Optional[dict] = None
    sources: Optional[dict] = None
class LoginIn(BaseModel):
    username: str
    password: str
class TwoFALoginIn(BaseModel):
    pending: str
    code: str
class PwChangeIn(BaseModel):
    old_password: str
    new_password: str
class CodeIn(BaseModel):
    code: str
class PwIn(BaseModel):
    password: str
class UserIn(BaseModel):
    username: Optional[str] = None
    label: Optional[str] = None
    password: Optional[str] = None
    admin: Optional[bool] = None
    groups: Optional[List[str]] = None
    reset_2fa: Optional[bool] = None
class GroupIn(BaseModel):
    id: Optional[str] = None
    label: Optional[str] = None
    folders: Optional[List[str]] = None

@app.get("/api/engines")
def engines(user: dict = Depends(auth.require_user)):
    cfg = load_config(); out = []
    for e in ORDER:
        en = cfg["engines"][e]
        label = LABELS[e] + (f" ({en['model']})" if e == "local" else "")
        out.append({"id": e, "label": label, "available": engine_available(cfg, e), "default": cfg["default_engine"] == e})
    return out

@app.get("/api/ollama/models")
def ollama_models(url: str = "", user: dict = Depends(auth.require_admin)):
    url = (url or "").strip().rstrip("/")
    if not url:
        return JSONResponse({"error": "Keine Ollama-URL angegeben."}, status_code=400)
    try:
        req = urllib.request.Request(url + "/api/tags")
        with urllib.request.urlopen(req, timeout=8) as r:
            data = json.load(r)
        models = sorted(m.get("name") for m in data.get("models", []) if m.get("name"))
        return {"models": models}
    except Exception as ex:
        return JSONResponse({"error": f"Ollama unter '{url}' nicht erreichbar: {ex}"}, status_code=400)

@app.post("/api/ask")
def ask(r: AskReq, user: dict = Depends(auth.require_user)):
    cfg = load_config(); e = r.engine or cfg["default_engine"]
    if e not in GENERATORS or not engine_available(cfg, e):
        return JSONResponse({"error": f"Engine '{e}' ist nicht verfuegbar (in Einstellungen aktivieren / Schluessel hinterlegen)."}, status_code=400)
    hits = retrieve(cfg, r.question, user)
    context = "\n\n".join(f"[Quelle: {h.payload['source']}]\n{h.payload['text']}" for h in hits)
    sources = sorted({h.payload["source"] for h in hits})
    try: answer = GENERATORS[e](cfg, r.question, context)
    except Exception as ex: return JSONResponse({"error": f"Fehler bei Engine '{LABELS.get(e, e)}': {ex}"}, status_code=500)
    return {"answer": answer, "sources": sources, "engine": LABELS[e]}

@app.post("/api/ask/stream")
def ask_stream(r: AskReq, user: dict = Depends(auth.require_user)):
    cfg = load_config(); e = r.engine or cfg["default_engine"]
    if e not in STREAMERS or not engine_available(cfg, e):
        def err():
            yield sse("error", {"error": f"Engine '{e}' ist nicht verfuegbar (in Einstellungen aktivieren / Schluessel hinterlegen)."})
        return StreamingResponse(err(), media_type="text/event-stream")
    hits = retrieve(cfg, r.question, user)
    context = "\n\n".join(f"[Quelle: {h.payload['source']}]\n{h.payload['text']}" for h in hits)
    sources = rank_sources(hits)
    messages = build_messages(r.history, r.question, context)
    def gen():
        yield sse("meta", {"sources": sources, "engine": LABELS[e]})
        try:
            for piece in STREAMERS[e](cfg, messages):
                yield sse("delta", {"text": piece})
        except Exception as ex:
            yield sse("error", {"error": f"Fehler bei Engine '{LABELS.get(e, e)}': {ex}"})
        yield sse("done", {})
    return StreamingResponse(gen(), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})

@app.get("/api/config")
def get_config(user: dict = Depends(auth.require_admin)):
    cfg = load_config(); raw = load_raw()
    out = {"default_engine": cfg["default_engine"], "retrieval": cfg["retrieval"], "engines": {}, "indexed_chunks": indexed_count()}
    for e in ORDER:
        en = cfg["engines"][e]; item = {"enabled": en.get("enabled", False), "model": en.get("model", ""), "label": LABELS[e]}
        if e == "local": item["ollama_url"] = en.get("ollama_url", "")
        else:
            item["has_key"] = bool(engine_key(cfg, e))
            # key_from_env anhand der Rohdaten (vor Env-Merge) bestimmen: kein Key in der Datei, aber in der Umgebung
            item["key_from_env"] = bool(not raw["engines"][e].get("api_key") and os.environ.get(ENV_KEYS.get(e, "")))
        out["engines"][e] = item
    s = cfg["sources"]["smb"]
    out["sources"] = {"smb": {"enabled": s.get("enabled", False), "host": s.get("host", ""), "share": s.get("share", ""),
                              "path": s.get("path", ""), "username": s.get("username", ""), "has_password": bool(s.get("password"))}}
    return out

@app.post("/api/config")
def set_config(c: ConfigIn, user: dict = Depends(auth.require_admin)):
    cfg = load_raw()
    if c.default_engine: cfg["default_engine"] = c.default_engine
    if c.retrieval: cfg["retrieval"].update(c.retrieval)
    if c.engines:
        for e, vals in c.engines.items():
            if e not in cfg["engines"]: continue
            tgt = cfg["engines"][e]
            if "enabled" in vals: tgt["enabled"] = bool(vals["enabled"])
            if vals.get("model"): tgt["model"] = vals["model"]
            if e == "local" and vals.get("ollama_url"): tgt["ollama_url"] = vals["ollama_url"]
            if vals.get("api_key"): tgt["api_key"] = vals["api_key"]
    if c.sources and "smb" in c.sources:
        smb = c.sources["smb"]; tgt = cfg["sources"]["smb"]
        if "enabled" in smb: tgt["enabled"] = bool(smb["enabled"])
        for f in ("host", "share", "path", "username"):
            if f in smb: tgt[f] = smb[f]
        if smb.get("password"): tgt["password"] = smb["password"]
    save_config(cfg)
    return {"ok": True}

@app.get("/api/status")
def api_status(user: dict = Depends(auth.require_user)):
    with _state_lock: st = dict(STATE)
    st["indexed_chunks"] = indexed_count(); return st

@app.post("/api/ingest")
def api_ingest(full: bool = False, user: dict = Depends(auth.require_admin)):
    # Asynchron: im Hintergrund-Thread starten, sofort zurueckkehren.
    # Fortschritt + Ergebnis holt das Frontend ueber /api/status.
    with _state_lock:
        if STATE["indexing"]:
            return JSONResponse({"error": "Es laeuft bereits eine Indexierung."}, status_code=409)
        STATE["indexing"] = True
        STATE["detail"] = "Start…"
        STATE["result"] = None
        STATE["error"] = None
    def worker():
        try:
            res = ingest_smb(full=full)
            with _state_lock: STATE["result"] = {"full": full, **res}
        except Exception as ex:
            with _state_lock: STATE["error"] = str(ex)
        finally:
            set_indexing(False)
    threading.Thread(target=worker, daemon=True).start()
    return {"ok": True, "started": True, "full": full}

@app.get("/api/documents")
def api_documents(user: dict = Depends(auth.require_user)):
    docs = list_documents(user)
    return {"documents": docs, "total_chunks": sum(d["chunks"] for d in docs)}

@app.get("/api/document")
def api_document(source: str, user: dict = Depends(auth.require_user)):
    # Nur tatsaechlich indizierte UND fuer den Nutzer freigegebene Quellen erlauben
    known = {d["source"] for d in list_documents(user)}
    if source not in known:
        return JSONResponse({"error": "Unbekannte Quelle."}, status_code=404)
    cfg = load_config(); s = cfg["sources"]["smb"]
    try:
        smbclient, root, base, sub = smb_connect(s)
        full = root + "\\" + source.replace("/", "\\")
        with smbclient.open_file(full, mode="rb", share_access="r") as fd: data = fd.read()
    except Exception as ex:
        return JSONResponse({"error": f"Datei nicht abrufbar: {ex}"}, status_code=502)
    name = source.rsplit("/", 1)[-1]
    ctype = mimetypes.guess_type(name)[0] or "application/octet-stream"
    # PDF/Text inline im Browser, Office-Dateien laden herunter
    disp = "inline" if ctype in ("application/pdf", "text/plain") else "attachment"
    headers = {"Content-Disposition": f"{disp}; filename*=UTF-8''{quote(name)}"}
    return Response(content=data, media_type=ctype, headers=headers)

@app.post("/api/upload")
async def api_upload(files: List[UploadFile] = File(...), folder: str = Form(""),
                     user: dict = Depends(auth.require_user)):
    cfg = load_config(); s = cfg["sources"]["smb"]
    if not s.get("share") or not s.get("username"):
        return JSONResponse({"error": "SMB-Quelle ist nicht konfiguriert (Share/Benutzer fehlen)."}, status_code=400)
    # Zielordner (relativ zur Share-Wurzel) bestimmen + Schreibrecht pruefen
    folder = (folder or "").strip().strip("/").replace("\\", "/")
    if not folder:
        if not user["admin"]:
            return JSONResponse({"error": "Bitte einen Zielordner waehlen."}, status_code=400)
        sub = (s.get("path") or "").strip().strip("\\/").replace("\\", "/")
        folder = sub  # Admin ohne Ordner -> konfigurierter Standardpfad
    if not src_allowed(user, folder):
        return JSONResponse({"error": "Kein Schreibrecht fuer diesen Ordner."}, status_code=403)
    try: smbclient, root, base, sub = smb_connect(s)
    except Exception as ex: return JSONResponse({"error": f"SMB-Verbindung fehlgeschlagen: {ex}"}, status_code=400)
    destdir = root + ("\\" + folder.replace("/", "\\") if folder else "")
    try:
        if folder: smbclient.makedirs(destdir, exist_ok=True)
    except Exception: pass
    results = []
    for uf in files:
        name = os.path.basename(uf.filename or "")
        if not name.lower().endswith(EXTS):
            results.append({"file": name, "error": "Dateityp nicht unterstuetzt"}); continue
        data = await uf.read()
        dest = destdir + "\\" + name
        try:
            with smbclient.open_file(dest, mode="wb") as fd: fd.write(data)
        except Exception as ex:
            results.append({"file": name, "error": f"Schreiben auf QNAP fehlgeschlagen: {ex}"}); continue
        rel = (folder + "/" if folder else "") + name
        try: sig = file_sig(smbclient.stat(dest))
        except Exception: sig = None
        try: results.append({"file": name, "chunks": index_file(rel, data, sig)})
        except Exception as ex: results.append({"file": name, "error": f"Indexierung fehlgeschlagen: {ex}"})
    return {"ok": True, "results": results, "indexed_chunks": indexed_count()}

def _warmup_reranker():
    try:
        rc = load_config().get("retrieval", {})
        if rc.get("rerank", True):
            get_reranker(rc.get("rerank_model", DEFAULT_RERANK)); print("[Warmup] Reranker bereit")
    except Exception as ex: print(f"[Warmup-Warn] {ex}")
# Warmup standardmäßig AUS: spart auf der 8-GB-VM ~1 GB RAM während des Indexierens
# (e5-large braucht den Platz). Der Reranker wird sonst beim ersten Chat lazy geladen
# (einmalig ~10 s Verzögerung). Mit RERANK_WARMUP=1 wieder aktivierbar.
if os.environ.get("RERANK_WARMUP", "0") == "1":
    threading.Thread(target=_warmup_reranker, daemon=True).start()

# ===================== Authentifizierung & Konto =====================
auth.bootstrap_admin()

@app.post("/api/login")
def login(r: LoginIn):
    uname = r.username.strip().lower()
    if auth.is_locked(uname):
        return JSONResponse({"error": "Zu viele Fehlversuche. Bitte kurz warten."}, status_code=429)
    store = auth.load_store(); u = store["users"].get(uname)
    if not u or not auth.verify_pw(r.password, u.get("pw_hash", "")):
        auth.note_fail(uname)
        return JSONResponse({"error": "Benutzername oder Passwort falsch."}, status_code=401)
    auth.note_success(uname)
    if u.get("totp_enabled"):
        return {"need_2fa": True, "pending": auth.create_pending(uname)}
    resp = JSONResponse({"ok": True}); auth.set_cookie(resp, auth.create_session(uname))
    return resp

@app.post("/api/login/2fa")
def login_2fa(r: TwoFALoginIn):
    uname = auth.consume_pending(r.pending)
    if not uname:
        return JSONResponse({"error": "Anmeldung abgelaufen, bitte erneut einloggen."}, status_code=400)
    store = auth.load_store(); u = store["users"].get(uname)
    if not u or not auth.totp_verify(u.get("totp_secret"), r.code):
        return JSONResponse({"error": "Code ungueltig."}, status_code=401)
    resp = JSONResponse({"ok": True}); auth.set_cookie(resp, auth.create_session(uname))
    return resp

@app.post("/api/logout")
def logout(request: Request):
    tok = request.cookies.get(auth.COOKIE_NAME)
    if tok: auth.destroy_session(tok)
    resp = JSONResponse({"ok": True}); auth.clear_cookie(resp); return resp

@app.get("/api/me")
def me(user: dict = Depends(auth.require_user)):
    return user

@app.post("/api/account/password")
def change_pw(r: PwChangeIn, user: dict = Depends(auth.require_user)):
    store = auth.load_store(); u = store["users"][user["username"]]
    if not auth.verify_pw(r.old_password, u.get("pw_hash", "")):
        return JSONResponse({"error": "Aktuelles Passwort falsch."}, status_code=401)
    if len(r.new_password) < 8:
        return JSONResponse({"error": "Neues Passwort muss mind. 8 Zeichen haben."}, status_code=400)
    u["pw_hash"] = auth.hash_pw(r.new_password); auth.save_store(store)
    return {"ok": True}

@app.post("/api/2fa/setup")
def twofa_setup(user: dict = Depends(auth.require_user)):
    store = auth.load_store(); u = store["users"][user["username"]]
    secret = auth.new_totp_secret()
    u["totp_secret"] = secret; u["totp_enabled"] = False  # erst nach Verify aktiv
    auth.save_store(store)
    uri = auth.totp_uri(secret, user["username"])
    return {"secret": secret, "otpauth_url": uri, "qr": auth.qr_png_base64(uri)}

@app.post("/api/2fa/enable")
def twofa_enable(r: CodeIn, user: dict = Depends(auth.require_user)):
    store = auth.load_store(); u = store["users"][user["username"]]
    if not auth.totp_verify(u.get("totp_secret"), r.code):
        return JSONResponse({"error": "Code ungueltig."}, status_code=401)
    u["totp_enabled"] = True; auth.save_store(store)
    return {"ok": True}

@app.post("/api/2fa/disable")
def twofa_disable(r: PwIn, user: dict = Depends(auth.require_user)):
    store = auth.load_store(); u = store["users"][user["username"]]
    if not auth.verify_pw(r.password, u.get("pw_hash", "")):
        return JSONResponse({"error": "Passwort falsch."}, status_code=401)
    u["totp_secret"] = None; u["totp_enabled"] = False; auth.save_store(store)
    return {"ok": True}

# ===================== Admin: Nutzer & Gruppen =====================
def _user_public(uname, u):
    return {"username": uname, "label": u.get("label", uname), "admin": bool(u.get("admin")),
            "groups": u.get("groups", []), "totp_enabled": bool(u.get("totp_enabled"))}
def _count_admins(store): return sum(1 for x in store["users"].values() if x.get("admin"))
def _slug(s):
    return re.sub(r"[^a-z0-9]+", "-", (s or "").strip().lower()).strip("-")

@app.get("/api/admin/users")
def admin_users(user: dict = Depends(auth.require_admin)):
    store = auth.load_store()
    return {"users": [_user_public(k, v) for k, v in sorted(store["users"].items())]}

@app.post("/api/admin/users")
def admin_user_create(r: UserIn, user: dict = Depends(auth.require_admin)):
    uname = (r.username or "").strip().lower()
    if not uname or not r.password:
        return JSONResponse({"error": "Benutzername und Passwort erforderlich."}, status_code=400)
    store = auth.load_store()
    if uname in store["users"]:
        return JSONResponse({"error": "Benutzer existiert bereits."}, status_code=409)
    store["users"][uname] = {"label": r.label or uname, "pw_hash": auth.hash_pw(r.password),
        "admin": bool(r.admin), "groups": r.groups or [], "totp_secret": None, "totp_enabled": False}
    auth.save_store(store)
    return {"ok": True}

@app.put("/api/admin/users/{uname}")
def admin_user_update(uname: str, r: UserIn, user: dict = Depends(auth.require_admin)):
    uname = uname.strip().lower()
    store = auth.load_store(); u = store["users"].get(uname)
    if not u: return JSONResponse({"error": "Unbekannter Benutzer."}, status_code=404)
    if r.label is not None: u["label"] = r.label
    if r.groups is not None: u["groups"] = r.groups
    if r.admin is not None:
        if not r.admin and u.get("admin") and _count_admins(store) <= 1:
            return JSONResponse({"error": "Der letzte Admin kann nicht entzogen werden."}, status_code=400)
        u["admin"] = bool(r.admin)
    if r.password: u["pw_hash"] = auth.hash_pw(r.password)
    if r.reset_2fa: u["totp_secret"] = None; u["totp_enabled"] = False
    auth.save_store(store)
    return {"ok": True}

@app.delete("/api/admin/users/{uname}")
def admin_user_delete(uname: str, user: dict = Depends(auth.require_admin)):
    uname = uname.strip().lower()
    store = auth.load_store(); u = store["users"].get(uname)
    if not u: return JSONResponse({"error": "Unbekannter Benutzer."}, status_code=404)
    if uname == user["username"]:
        return JSONResponse({"error": "Du kannst dich nicht selbst loeschen."}, status_code=400)
    if u.get("admin") and _count_admins(store) <= 1:
        return JSONResponse({"error": "Der letzte Admin kann nicht geloescht werden."}, status_code=400)
    del store["users"][uname]; auth.save_store(store)
    return {"ok": True}

@app.get("/api/admin/groups")
def admin_groups(user: dict = Depends(auth.require_admin)):
    store = auth.load_store()
    return {"groups": [{"id": k, "label": v.get("label", k), "folders": v.get("folders", [])}
                       for k, v in sorted(store["groups"].items())]}

@app.post("/api/admin/groups")
def admin_group_create(r: GroupIn, user: dict = Depends(auth.require_admin)):
    gid = _slug(r.id or r.label)
    if not gid: return JSONResponse({"error": "Name erforderlich."}, status_code=400)
    store = auth.load_store()
    if gid in store["groups"]:
        return JSONResponse({"error": "Gruppe existiert bereits."}, status_code=409)
    store["groups"][gid] = {"label": r.label or gid, "folders": r.folders or []}
    auth.save_store(store)
    return {"ok": True, "id": gid}

@app.put("/api/admin/groups/{gid}")
def admin_group_update(gid: str, r: GroupIn, user: dict = Depends(auth.require_admin)):
    store = auth.load_store(); g = store["groups"].get(gid)
    if not g: return JSONResponse({"error": "Unbekannte Gruppe."}, status_code=404)
    if r.label is not None: g["label"] = r.label
    if r.folders is not None: g["folders"] = r.folders
    auth.save_store(store)
    return {"ok": True}

@app.delete("/api/admin/groups/{gid}")
def admin_group_delete(gid: str, user: dict = Depends(auth.require_admin)):
    store = auth.load_store()
    if gid not in store["groups"]:
        return JSONResponse({"error": "Unbekannte Gruppe."}, status_code=404)
    del store["groups"][gid]
    for u in store["users"].values():
        if gid in u.get("groups", []): u["groups"] = [x for x in u["groups"] if x != gid]
    auth.save_store(store)
    return {"ok": True}

@app.get("/api/admin/folders")
def admin_folders(user: dict = Depends(auth.require_admin)):
    return {"folders": all_prefixes()}

app.mount("/", StaticFiles(directory="static", html=True), name="static")
